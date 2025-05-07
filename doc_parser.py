# doc_parser.py
import io
from PyPDF2 import PdfReader
from pptx import Presentation

def extract_text_from_pdf(file_bytes):
    """Extracts text from a PDF file given as bytes."""
    try:
        pdf_file = io.BytesIO(file_bytes)
        reader = PdfReader(pdf_file)
        text = ""
        for page in reader.pages:
            text += page.extract_text() or ""
        return text
    except Exception as e:
        print(f"Error parsing PDF: {e}")
        return ""

def extract_text_from_pptx(file_bytes):
    """Extracts text from a PowerPoint file given as bytes."""
    try:
        pptx_file = io.BytesIO(file_bytes)
        prs = Presentation(pptx_file)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        print(f"Error parsing PPTX: {e}")
        return ""

def extract_text_from_uploaded_files(uploaded_files):
    """
    Extracts text from a list of Streamlit UploadedFile objects.
    Supports PDF and PPTX.
    """
    combined_text = ""
    if not uploaded_files:
        return ""
        
    for uploaded_file in uploaded_files:
        file_bytes = uploaded_file.getvalue()
        file_name = uploaded_file.name.lower()
        
        if file_name.endswith(".pdf"):
            print(f"Parsing PDF: {uploaded_file.name}")
            combined_text += extract_text_from_pdf(file_bytes) + "\n\n"
        elif file_name.endswith(".pptx"):
            print(f"Parsing PPTX: {uploaded_file.name}")
            combined_text += extract_text_from_pptx(file_bytes) + "\n\n"
        else:
            print(f"Unsupported file type: {uploaded_file.name}. Skipping.")
            
    return combined_text.strip()